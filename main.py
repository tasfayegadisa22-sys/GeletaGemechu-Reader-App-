import os
from flask import Flask, request, send_file, render_template, make_response
from gtts import gTTS
import PyPDF2
import docx
import pptx
import zipfile
from bs4 import BeautifulSoup
import edge_tts
import asyncio
import openpyxl
import pytesseract
from PIL import Image

app = Flask(__name__)

async def _generate_edge_tts(text, voice, filename):
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(filename)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/read', methods=['POST'])
def read_text():
    text = request.form.get('text', '')
    lang = request.form.get('lang', 'am') 
    file = request.files.get('file')

    if file:
        filename = file.filename.lower()
        try:
            if filename.endswith('.txt'):
                text = file.read().decode('utf-8')
            elif filename.endswith('.pdf'):
                reader = PyPDF2.PdfReader(file)
                text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
            elif filename.endswith('.docx'):
                doc = docx.Document(file)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif filename.endswith('.pptx'):
                prs = pptx.Presentation(file)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = "\n".join(text_runs)
            elif filename.endswith('.epub'):
                text_runs = []
                with zipfile.ZipFile(file) as archive:
                    for item in archive.namelist():
                        if item.endswith('.html') or item.endswith('.xhtml'):
                            content = archive.read(item)
                            soup = BeautifulSoup(content, 'html.parser')
                            text_runs.append(soup.get_text())
                text = "\n".join(text_runs)
            elif filename.endswith('.xlsx'):
                # የኤክሴል ፋይል ማንበቢያ
                wb = openpyxl.load_workbook(file)
                text_runs = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        for cell in row:
                            if cell:
                                text_runs.append(str(cell))
                text = " ".join(text_runs)
            elif filename.endswith(('.png', '.jpg', '.jpeg')):
                # የፎቶ/ምስል ማንበቢያ (OCR)
                image = Image.open(file)
                # አማርኛ እና እንግሊዝኛን አጣምሮ እንዲያነብ 'amh+eng' እንለዋለን
                text = pytesseract.image_to_string(image, lang='amh+eng')
            else:
                return "ይህ የፋይል አይነት አይደገፍም።", 400
        except Exception as e:
            return f"ፋይሉን ማንበብ አልተቻለም: {str(e)}", 500
    
    if not text.strip():
        return "ምንም ጽሁፍ አልተገኘም ወይም ከፎቶው ላይ ጽሁፍ መለየት አልተቻለም", 400

    word_count = len(text.split())
    audio_minutes = max(1, round(word_count / 130))
    
    try:
        voice = 'en-US-AriaNeural' if lang == 'en' else 'am-ET-AmehaNeural'
        asyncio.run(_generate_edge_tts(text, voice, "speech.mp3"))
        
        response = make_response(send_file("speech.mp3", mimetype="audio/mpeg"))
        response.headers['X-Word-Count'] = str(word_count)
        response.headers['X-Audio-Minutes'] = str(audio_minutes)
        return response

    except Exception as edge_error:
        try:
            tts = gTTS(text=text, lang=lang, slow=False)
            tts.save("speech.mp3")
            response = make_response(send_file("speech.mp3", mimetype="audio/mpeg"))
            response.headers['X-Word-Count'] = str(word_count)
            response.headers['X-Audio-Minutes'] = str(audio_minutes)
            return response
        except Exception as gtts_error:
            return f"የድምጽ ማመንጨት አልተቻለም: {str(edge_error)}", 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 10000))
    app.run(host='0.0.0.0', port=port)
